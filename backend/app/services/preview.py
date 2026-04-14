"""Generate HTML previews for office documents (Excel, PowerPoint)."""

import asyncio
import logging
import os
from pathlib import Path

logger = logging.getLogger(__name__)

_HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="utf-8">
<style>
  html, body {{ margin: 0; padding: 0; overflow: hidden; height: 100vh; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; color: #222; background: #fff; }}
  @media (prefers-color-scheme: dark) {{ body {{ background: #1a1a1a; color: #e0e0e0; }} }}
  .table-wrap {{ overflow: auto; height: 100vh; }}
  table {{ border-collapse: collapse; width: max-content; min-width: 100%; font-size: 13px; }}
  th, td {{ border: 1px solid #ccc; padding: 6px 10px; text-align: left; white-space: nowrap; }}
  @media (prefers-color-scheme: dark) {{ th, td {{ border-color: #444; }} }}
  th {{ background: #f0f0f0; font-weight: 600; }}
  @media (prefers-color-scheme: dark) {{ th {{ background: #2a2a2a; }} }}
  h2 {{ font-size: 16px; margin: 24px 0 8px; padding-bottom: 4px; border-bottom: 2px solid #ddd; }}
  @media (prefers-color-scheme: dark) {{ h2 {{ border-color: #444; }} }}
  .slide {{ border: 1px solid #ccc; border-radius: 8px; padding: 20px; margin-bottom: 20px; }}
  @media (prefers-color-scheme: dark) {{ .slide {{ border-color: #444; }} }}
  .slide-num {{ font-size: 12px; color: #888; margin-bottom: 8px; }}
  .slide p {{ margin: 4px 0; }}
</style>
</head>
<body>{content}</body>
</html>"""


async def render_preview_html(path: str, file_type: str) -> str:
    """Render an office document as HTML."""
    ft = file_type.lower().lstrip(".")
    if ft == "xls":
        converted = await _convert_for_preview(path, "xls", "xlsx")
        if converted:
            return await asyncio.get_event_loop().run_in_executor(None, _render_excel, converted)
        return _HTML_TEMPLATE.format(content="<p>プレビューの生成に失敗しました。</p>")
    elif ft == "xlsx":
        return await asyncio.get_event_loop().run_in_executor(None, _render_excel, path)
    elif ft == "ppt":
        converted = await _convert_for_preview(path, "ppt", "pptx")
        if converted:
            return await asyncio.get_event_loop().run_in_executor(None, _render_pptx, converted)
        return _HTML_TEMPLATE.format(content="<p>プレビューの生成に失敗しました。</p>")
    elif ft == "pptx":
        return await asyncio.get_event_loop().run_in_executor(None, _render_pptx, path)
    elif ft in ("csv", "tsv"):
        return await asyncio.get_event_loop().run_in_executor(None, _render_csv, path, ft)
    elif ft == "doc":
        converted = await _convert_for_preview(path, "doc", "docx")
        if converted:
            return await asyncio.get_event_loop().run_in_executor(None, _render_docx, converted)
        return _HTML_TEMPLATE.format(content="<p>プレビューの生成に失敗しました。</p>")
    elif ft == "docx":
        return await asyncio.get_event_loop().run_in_executor(None, _render_docx, path)
    else:
        return _HTML_TEMPLATE.format(content="<p>このファイル形式のプレビューには対応していません。</p>")


async def _convert_for_preview(path: str, src_ext: str, dst_ext: str) -> str | None:
    """Convert a legacy Office file for preview using LibreOffice headless."""
    import tempfile
    import uuid as _uuid

    tmpdir = tempfile.mkdtemp()
    try:
        symlink = os.path.join(tmpdir, f"{_uuid.uuid4().hex}.{src_ext}")
        os.symlink(os.path.abspath(path), symlink)
        user_inst = f"file://{tmpdir}/lo_profile"
        proc = await asyncio.create_subprocess_exec(
            "soffice", "--headless", "--norestore",
            f"-env:UserInstallation={user_inst}",
            "--convert-to", dst_ext, "--outdir", tmpdir, symlink,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        )
        await asyncio.wait_for(proc.communicate(), timeout=120)
        converted = symlink.rsplit(".", 1)[0] + f".{dst_ext}"
        if os.path.exists(converted):
            return converted
        logger.error("LibreOffice preview conversion failed: %s -> %s", path, dst_ext)
        return None
    except Exception as e:
        logger.error("LibreOffice preview conversion error for %s: %s", path, e)
        return None


def _render_excel(path: str) -> str:
    from io import BytesIO
    from openpyxl import load_workbook

    # openpyxl checks file extension — TUS uploads have no extension,
    # so read into BytesIO to bypass the filename check.
    with open(path, "rb") as f:
        buf = BytesIO(f.read())
    wb = load_workbook(buf, read_only=True, data_only=True)
    parts: list[str] = []

    multi_sheet = len(wb.sheetnames) > 1
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append('<div class="table-wrap">')
        if multi_sheet:
            parts.append(f"<h2>{_esc(sheet_name)}</h2>")
        parts.append("<table>")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(c) if c is not None else "" for c in row]
            if not any(c for c in cells):
                continue
            tag = "th" if i == 0 else "td"
            parts.append("<tr>" + "".join(f"<{tag}>{_esc(c)}</{tag}>" for c in cells) + "</tr>")
        parts.append("</table>\n</div>")

    wb.close()
    return _HTML_TEMPLATE.format(content="\n".join(parts))


def _render_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'<div class="slide">')
        parts.append(f'<div class="slide-num">スライド {i}</div>')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(f"<p>{_esc(text)}</p>")
            if shape.has_table:
                parts.append("<table>")
                for ri, row in enumerate(shape.table.rows):
                    tag = "th" if ri == 0 else "td"
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append("<tr>" + "".join(f"<{tag}>{_esc(c)}</{tag}>" for c in cells) + "</tr>")
                parts.append("</table>")
        parts.append("</div>")

    return _HTML_TEMPLATE.format(content="\n".join(parts))


def _render_csv(path: str, ft: str) -> str:
    import csv

    delimiter = "\t" if ft == "tsv" else ","
    parts: list[str] = ['<div class="table-wrap"><table>']
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        for i, row in enumerate(reader):
            if not any(c.strip() for c in row):
                continue
            tag = "th" if i == 0 else "td"
            parts.append("<tr>" + "".join(f"<{tag}>{_esc(c)}</{tag}>" for c in row) + "</tr>")
    parts.append("</table></div>")
    return _HTML_TEMPLATE.format(content="\n".join(parts))


def _render_docx(path: str) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(f"<p>{_esc(text)}</p>")
    for table in doc.tables:
        parts.append("<table>")
        for ri, row in enumerate(table.rows):
            tag = "th" if ri == 0 else "td"
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("<tr>" + "".join(f"<{tag}>{_esc(c)}</{tag}>" for c in cells) + "</tr>")
        parts.append("</table>")
    return _HTML_TEMPLATE.format(content="\n".join(parts))


def _esc(s: str) -> str:
    """Escape HTML special characters."""
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")
