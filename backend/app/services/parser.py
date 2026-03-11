import asyncio
import csv
import io
import logging
import os
from pathlib import Path

logger = logging.getLogger(__name__)

# Image extensions that should be processed via OCR
_IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "gif", "bmp", "tiff", "tif", "webp"}

# OCR server URL (host-side Surya server)
_OCR_SERVER_URL = os.environ.get("OCR_SERVER_URL", "http://host.docker.internal:8090")


async def parse_file(path: str, file_type: str) -> str:
    """Extract text content from a file based on its type."""
    file_type = file_type.lower().lstrip(".")

    if file_type in ("md", "txt", "text", "markdown"):
        return await _parse_text(path)
    elif file_type == "pdf":
        return await _parse_pdf(path)
    elif file_type in ("docx", "doc"):
        return await _parse_docx(path)
    elif file_type in ("xlsx", "xls"):
        return await _parse_excel(path)
    elif file_type in ("csv", "tsv"):
        return await _parse_csv(path, delimiter="\t" if file_type == "tsv" else ",")
    elif file_type in ("htm", "html"):
        return await _parse_html(path)
    elif file_type == "pptx":
        return await _parse_pptx(path)
    elif file_type in _IMAGE_EXTENSIONS:
        return await _parse_image_ocr(path)
    else:
        # Attempt plain-text read as fallback
        return await _parse_text(path)


async def _parse_text(path: str) -> str:
    """Read a plain text / markdown file."""

    def _read() -> str:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    return await asyncio.get_event_loop().run_in_executor(None, _read)


async def _parse_pdf(path: str) -> str:
    """Extract text from a PDF using PyMuPDF (fitz).

    If the PDF contains very little text (likely scanned/image-based),
    falls back to OCR via the external OCR server.
    """

    def _extract() -> str:
        import fitz  # PyMuPDF

        text_parts: list[str] = []
        with fitz.open(path) as doc:
            for page in doc:
                text_parts.append(page.get_text())
        return "\n".join(text_parts)

    text = await asyncio.get_event_loop().run_in_executor(None, _extract)

    # If extracted text is too short, it's likely a scanned PDF — try OCR
    stripped = text.strip()
    if len(stripped) < 50:
        logger.info("PDF has little text (%d chars), attempting OCR: %s", len(stripped), path)
        ocr_text = await _call_ocr_server(path)
        if ocr_text:
            return ocr_text

    return text


async def _parse_docx(path: str) -> str:
    """Extract text from a DOCX file."""

    def _extract() -> str:
        from docx import Document as DocxDocument

        doc = DocxDocument(path)
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def _parse_excel(path: str) -> str:
    """Extract text from an Excel file (.xlsx/.xls)."""

    def _extract() -> str:
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                # Skip completely empty rows
                if any(c for c in cells):
                    parts.append(" | ".join(cells))
            parts.append("")  # blank line between sheets

        wb.close()
        return "\n".join(parts)

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def _parse_csv(path: str, delimiter: str = ",") -> str:
    """Extract text from a CSV/TSV file."""

    def _extract() -> str:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f, delimiter=delimiter)
            lines: list[str] = []
            for row in reader:
                if any(c.strip() for c in row):
                    lines.append(" | ".join(row))
        return "\n".join(lines)

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def _parse_html(path: str) -> str:
    """Extract text from an HTML file."""

    def _extract() -> str:
        from bs4 import BeautifulSoup

        with open(path, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f, "html.parser")

        # Remove script and style elements
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        text = soup.get_text(separator="\n")
        # Clean up excessive blank lines
        lines = [line.strip() for line in text.splitlines()]
        return "\n".join(line for line in lines if line)

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def _parse_pptx(path: str) -> str:
    """Extract text from a PowerPoint file (.pptx)."""

    def _extract() -> str:
        from pptx import Presentation

        prs = Presentation(path)
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                parts.append(f"## スライド {i}")
                parts.extend(slide_texts)
                parts.append("")

        return "\n".join(parts)

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def _parse_image_ocr(path: str) -> str:
    """Extract text from an image file via the OCR server."""
    text = await _call_ocr_server(path)
    if not text:
        logger.warning("OCR returned no text for image: %s", path)
        return ""
    return text


async def _call_ocr_server(path: str) -> str:
    """Send a file to the OCR server and return extracted text."""
    import httpx

    url = f"{_OCR_SERVER_URL}/ocr"
    filename = Path(path).name

    try:
        async with httpx.AsyncClient(timeout=600.0) as client:
            with open(path, "rb") as f:
                resp = await client.post(url, files={"file": (filename, f)})
            resp.raise_for_status()
            return resp.json().get("text", "")
    except httpx.ConnectError:
        logger.warning("OCR server not available at %s", _OCR_SERVER_URL)
        return ""
    except Exception as e:
        logger.error("OCR server error: %s", e)
        return ""


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    """Split text into overlapping chunks by character count.

    Attempts to break on paragraph / sentence boundaries when possible.
    """
    if not text or not text.strip():
        return []

    chunks: list[str] = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + chunk_size, text_len)

        # If we're not at the very end, try to find a good break point
        if end < text_len:
            # Look for paragraph break first
            para_break = text.rfind("\n\n", start, end)
            if para_break > start + chunk_size // 2:
                end = para_break + 2  # include the double newline
            else:
                # Look for single newline
                newline = text.rfind("\n", start + chunk_size // 2, end)
                if newline > start:
                    end = newline + 1
                else:
                    # Look for sentence end
                    for sep in (". ", "! ", "? "):
                        pos = text.rfind(sep, start + chunk_size // 2, end)
                        if pos > start:
                            end = pos + len(sep)
                            break

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        # Move start forward, applying overlap
        start = max(end - overlap, start + 1)
        if start >= text_len:
            break

    return chunks
